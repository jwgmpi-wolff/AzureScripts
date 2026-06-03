"""
Document Processing Pipeline for FinOps Knowledge Base
Handles PPTX and other document formats, chunking, and embedding
"""
import os
import re
import json
from pathlib import Path
from typing import List, Dict, Any, Tuple
from datetime import datetime
import logging

# For PPTX processing
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed. Install with: pip install python-pptx")

# For PDF processing
try:
    import PyPDF2
except ImportError:
    print("PyPDF2 not installed. Install with: pip install PyPDF2")

from langchain_text_splitters import RecursiveCharacterTextSplitter

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process various document formats into structured chunks"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ".", " ", ""]
        )
    
    def process_pptx(self, file_path: str, category: str) -> List[Dict[str, str]]:
        """Extract text from PowerPoint file"""
        logger.info(f"Processing PPTX: {file_path}")
        chunks = []
        
        try:
            prs = Presentation(file_path)
            full_text = []
            slide_contents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"## Slide {slide_num}\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip() != f"## Slide {slide_num}\n":
                    slide_contents.append(slide_text)
                    full_text.append(slide_text)
            
            # Split content into chunks
            combined_text = "\n\n".join(full_text)
            text_chunks = self.splitter.split_text(combined_text)
            
            filename = Path(file_path).stem
            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "id": f"{category}_{filename}_{idx}",
                    "content": chunk,
                    "source": f"{filename}.pptx",
                    "category": category,
                    "title": f"{filename} - Part {idx + 1}",
                    "metadata": json.dumps({
                        "file_type": "pptx",
                        "processed_date": datetime.now().isoformat(),
                        "chunk_index": idx,
                        "original_file": file_path
                    })
                })
            
            logger.info(f"✓ Extracted {len(chunks)} chunks from {filename}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return []
    
    def process_pdf(self, file_path: str, category: str) -> List[Dict[str, str]]:
        """Extract text from PDF file"""
        logger.info(f"Processing PDF: {file_path}")
        chunks = []
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                full_text = []
                
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        full_text.append(f"## Page {page_num}\n{text}")
            
            combined_text = "\n\n".join(full_text)
            text_chunks = self.splitter.split_text(combined_text)
            
            filename = Path(file_path).stem
            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "id": f"{category}_{filename}_{idx}",
                    "content": chunk,
                    "source": f"{filename}.pdf",
                    "category": category,
                    "title": f"{filename} - Part {idx + 1}",
                    "metadata": json.dumps({
                        "file_type": "pdf",
                        "processed_date": datetime.now().isoformat(),
                        "chunk_index": idx,
                        "original_file": file_path
                    })
                })
            
            logger.info(f"✓ Extracted {len(chunks)} chunks from {filename}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []
    
    def process_markdown(self, file_path: str, category: str) -> List[Dict[str, str]]:
        """Process Markdown file"""
        logger.info(f"Processing Markdown: {file_path}")
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            text_chunks = self.splitter.split_text(content)
            filename = Path(file_path).stem
            
            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "id": f"{category}_{filename}_{idx}",
                    "content": chunk,
                    "source": f"{filename}.md",
                    "category": category,
                    "title": f"{filename} - Part {idx + 1}",
                    "metadata": json.dumps({
                        "file_type": "markdown",
                        "processed_date": datetime.now().isoformat(),
                        "chunk_index": idx,
                        "original_file": file_path
                    })
                })
            
            logger.info(f"✓ Extracted {len(chunks)} chunks from {filename}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing Markdown {file_path}: {e}")
            return []
    
    def process_text(self, file_path: str, category: str) -> List[Dict[str, str]]:
        """Process plain text file"""
        logger.info(f"Processing Text: {file_path}")
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            text_chunks = self.splitter.split_text(content)
            filename = Path(file_path).stem
            
            for idx, chunk in enumerate(text_chunks):
                chunks.append({
                    "id": f"{category}_{filename}_{idx}",
                    "content": chunk,
                    "source": f"{filename}.txt",
                    "category": category,
                    "title": f"{filename} - Part {idx + 1}",
                    "metadata": json.dumps({
                        "file_type": "text",
                        "processed_date": datetime.now().isoformat(),
                        "chunk_index": idx,
                        "original_file": file_path
                    })
                })
            
            logger.info(f"✓ Extracted {len(chunks)} chunks from {filename}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing Text {file_path}: {e}")
            return []
    
    def process_directory(self, kb_dir: str) -> List[Dict[str, str]]:
        """Process all documents in knowledge base directory structure"""
        all_chunks = []
        kb_path = Path(kb_dir)
        
        # Map directories to categories
        categories = {
            "aks": "AKS Cost Optimization",
            "vm": "Virtual Machine Optimization",
            "sql": "SQL Database & Cost",
            "postgres": "PostgreSQL & Open Source",
            "governance": "Governance & Policy",
            "tagging": "Tagging & Metadata",
            "cost-optimization": "Cost Optimization",
            "reserved-instances": "Reserved Instances & Savings Plans"
        }
        
        for category, display_name in categories.items():
            cat_path = kb_path / category
            if cat_path.exists():
                logger.info(f"\n📁 Processing category: {display_name}")
                for file_path in sorted(cat_path.glob("*")):
                    if file_path.is_file():
                        suffix = file_path.suffix.lower()
                        try:
                            if suffix == ".pptx":
                                chunks = self.process_pptx(str(file_path), category)
                            elif suffix == ".pdf":
                                chunks = self.process_pdf(str(file_path), category)
                            elif suffix == ".md":
                                chunks = self.process_markdown(str(file_path), category)
                            elif suffix == ".txt":
                                chunks = self.process_text(str(file_path), category)
                            else:
                                logger.warning(f"⚠️ Unsupported file type: {suffix}")
                                continue
                            
                            all_chunks.extend(chunks)
                        except Exception as e:
                            logger.error(f"Failed to process {file_path}: {e}")
        
        logger.info(f"\n✓ Total chunks processed: {len(all_chunks)}")
        return all_chunks


def save_chunks_to_jsonl(chunks: List[Dict[str, str]], output_file: str):
    """Save chunks to JSONL format for batch processing"""
    with open(output_file, 'w', encoding='utf-8') as f:
        for chunk in chunks:
            f.write(json.dumps(chunk, ensure_ascii=False) + '\n')
    logger.info(f"✓ Saved {len(chunks)} chunks to {output_file}")


if __name__ == "__main__":
    from config import AzureConfig
    
    config = AzureConfig()
    processor = DocumentProcessor(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP
    )
    
    # Process all documents
    chunks = processor.process_directory(config.KB_DIR)
    
    # Save to JSONL
    save_chunks_to_jsonl(chunks, "chunks.jsonl")
